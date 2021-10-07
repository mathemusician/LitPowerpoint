# import the streamlit library
import streamlit as st
import sys
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.util import Inches
from itertools import zip_longest
from math import ceil


def show_cat():
    pics = {
        "Cat": "https://cdn.pixabay.com/photo/2016/09/24/22/20/cat-1692702_960_720.jpg",
        "Puppy": "https://cdn.pixabay.com/photo/2019/03/15/19/19/puppy-4057786_960_720.jpg",
    }

    pic = st.selectbox("", list(pics.keys()), 0)
    st.image(pics[pic], use_column_width=True, caption=pics[pic])


def grouper(n, iterable, padvalue=""):
    """
    grouper(3, 'abcdefg', 'x') -->
        ('a','b','c')
        ('d','e','f')
        ('g','x','x')
    """
    return zip_longest(*[iter(iterable)] * n, fillvalue=padvalue)


def split_lyrics_by_number(lyrics_and_sections, split_number):

    groups = grouper(int(split_number), lyrics_and_sections)
    num_iters = ceil(len(lyrics_and_sections) / int(split_number))
    lyrics_and_sections = ["\n".join(next(groups)) for i in range(num_iters)]

    # get rid of trailing newlines
    lyrics_and_sections[-1] = lyrics_and_sections[-1].strip()

    return lyrics_and_sections


def split_by_font_size(lyrics_and_sections):
    """
    Parse lyrics by splitting them into separate lines.

    Parameters
    ----------
    lyrics_and_sections : list of str
        Lyrics to be parsed. Each element is a line in the song.
        
    Returns
    -------
    parsed_lyrics : list of str
        Parsed lyrics, split by line. Each element is a line in the song. 

    Example
    -------
    >>> split_by_font_size(["Some words are longer than 50 characters, so we need to split them."])
    ['Some words are longer than 50 characters, so we need to', 'split them.']
    """
    parsed_lyrics = []

    for lyric in lyrics_and_sections:

        if len(lyric) > 50:
            # just split it in half
            lyric_list = lyric.split(" ")
            length_of_lyric = len(lyric_list)

            if "\n" in lyric:
                length_of_first_part = lyric.index("\n")
                lyric = lyric.replace("\n", "")
                parsed_lyrics.append(lyric[:length_of_first_part])
                parsed_lyrics.append(lyric[length_of_first_part:])

            else:
                length_of_first_part = int(length_of_lyric / 2)
                first_part = lyric_list[:length_of_first_part]
                second_part = lyric_list[length_of_first_part:]
                parsed_lyrics.append(" ".join(first_part))
                parsed_lyrics.append(" ".join(second_part))

        else:
            parsed_lyrics.append(lyric)

    return parsed_lyrics


def get_rid_of_informative_lyrics(lyrics_and_sections):

    parsed_lyrics = []

    for lyric in lyrics_and_sections:
        if ("[" in lyric) or ("]" in lyric):
            pass
        else:
            parsed_lyrics.append(lyric)

    return parsed_lyrics


def make_powerpoint(font_size, split_lyrics_by, powerpoint_name, lyrics_and_sections):
    # split lyrics by ones
    lyrics_and_sections = lyrics_and_sections.split("\n")

    # get rid of informative lyrics
    lyrics_and_sections = get_rid_of_informative_lyrics(lyrics_and_sections)

    # split lyrics using font sizes
    lyrics_and_sections = split_by_font_size(lyrics_and_sections)

    # split lyrics by number
    lyrics_and_sections = split_lyrics_by_number(lyrics_and_sections, split_lyrics_by)

    # Make presentation
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]

    # make slides
    for lyric in lyrics_and_sections:

        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = lyric

        # make background black
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.theme_color = 1

        # Center text
        title.left = Inches(0)
        title.width = Inches(10)
        title.height = 1470025
        if split_lyrics_by == "1":
            title.top = Inches(3)
        elif split_lyrics_by == "2":
            title.top = Inches(3)
        elif split_lyrics_by == "3":
            title.top = Inches(2.5)
        elif split_lyrics_by == "4":
            title.top = Inches(2)

        # format individual lines of words
        for para_id in range(len(slide.shapes.title.text_frame.paragraphs)):
            title.text_frame.paragraphs[para_id].font.bold = True
            title.text_frame.paragraphs[para_id].font.name = "Helvetica"
            title.text_frame.paragraphs[para_id].font.size = Pt(
                font_size
            )  # can be 58 or 44
            title.text_frame.paragraphs[para_id].font.color.rgb = RGBColor(
                255, 255, 255
            )

    prs.save(f"{powerpoint_name}.pptx")


def main():
    st.title("Text to Powerpoint")

    font_size = st.number_input("Font size", value=50)

    split_lyrics_by = st.selectbox("Split Lyrics:", [str(i) for i in range(1, 5)], 1)

    powerpoint_name = st.text_input("Powerpoint Name:", value="Powerpoint Name")

    lyrics_and_sections = st.text_area("Lyrics:", value="Example Lyrics\nYay!")

    make_powerpoint(font_size, split_lyrics_by, powerpoint_name, lyrics_and_sections)

    with open(f"{powerpoint_name}.pptx", "rb") as file:
        btn = st.download_button(
            label="Download Powerpoint",
            data=file,
            file_name=f"{powerpoint_name}.pptx",
        )

    show_cat()


if __name__ == "__main__":
    main()
